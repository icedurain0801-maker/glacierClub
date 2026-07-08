import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


ROOT = Path.cwd()
MANIFEST = ROOT / ".codex-temp" / "next-work-suggestions" / "manifest.json"
OUTPUT = ROOT / "next-work-suggestions.pptx"


def main():
    slides = json.loads(MANIFEST.read_text(encoding="utf-8"))
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for item in slides:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(Path(item["image"])),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
